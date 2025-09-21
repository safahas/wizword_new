import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

ROOT = os.path.dirname(os.path.dirname(__file__))

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = item
        p.level = 0

def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    path = image_path if os.path.isabs(image_path) else os.path.join(ROOT, image_path)
    if os.path.exists(path):
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        slide.shapes.add_picture(path, left, top, width=width)
    else:
        # Fallback text when image missing
        tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
        tf = tx.text_frame
        tf.text = f"Image not found: {image_path}"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def main(output_path="docs/wizword_pitch.pptx"):
    prs = Presentation()
    add_title_slide(prs, "WizWord — Product Pitch", "AI‑assisted word game with shareable moments")

    add_bullets_slide(prs, "Problem & Opportunity", [
        "Word games often feel shallow or slow",
        "Need quick, meaningful sessions with progress",
    ])

    add_bullets_slide(prs, "What WizWord Does", [
        "AI‑assisted yes/no & hints keep momentum",
        "5‑minute Beat sprint with dynamic scoring",
        "Beautiful share cards with QR",
        "Optional Personal mode (privacy aware)",
    ])

    # How to Play (Brief) slide for pitch deck
    add_bullets_slide(prs, "How to Play (Brief)", [
        "Choose mode: Fun, Wiz, or Beat (4‑minute sprint)",
        "Pick a category or 'Any'",
        "Actions: yes/no (−1), hints (−10 up to 3), guess anytime",
        "Scoring: correct +20×length; wrong −10; skip −10",
        "SEI efficiency powers leaderboards; Personal gives profile‑aware hints",
    ])

    add_image_slide(prs, "Account Options", "assets/ui/account_options.png")
    add_image_slide(prs, "How to Play", "assets/ui/how_to_play.png")
    add_image_slide(prs, "Score Trend", "score_trend.png")
    add_image_slide(prs, "Score Distribution", "score_distribution.png")
    add_image_slide(prs, "Share Card Template", "assets/share_card_template.png")

    add_bullets_slide(prs, "Key Features", [
        "Fast play on web/mobile",
        "SEI rewards speed + accuracy",
        "Share cards with QR",
        "Offline fallback (no API key required)",
    ])

    add_bullets_slide(prs, "Setup (1‑minute)", [
        "python -m venv venv",
        "pip install -r requirements.txt",
        "streamlit run streamlit_app.py",
    ])

    out = os.path.join(ROOT, output_path)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()


